# service/app/tools/file_extractors/pptx_extractor.py
from __future__ import annotations

from typing import Optional

from ..vision_tool import GeminiVision, load_prompt

PROMPT_PATH = "packages/prompts/vision_extract_rich.md"


def extract_pptx(content: bytes, *, vision: Optional[GeminiVision], limits: dict) -> str:
    try:
        import io
        from pptx import Presentation
    except Exception:
        return ""

    prompt = load_prompt(PROMPT_PATH)
    max_imgs = int(limits.get("PPTX_VISION_MAX_IMAGES", 10))

    prs = Presentation(io.BytesIO(content))
    parts = []

    for si, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            try:
                if hasattr(shape, "text") and shape.text:
                    t = shape.text.strip()
                    if t:
                        texts.append(t)
            except Exception:
                pass
        parts.append(f"\n\n--- SLIDE {si} TEXT ---\n" + ("\n".join(texts) if texts else ""))

        if vision and vision.enabled() and max_imgs > 0:
            img_count = 0
            for shape in slide.shapes:
                if img_count >= max_imgs:
                    break
                try:
                    if shape.shape_type == 13:  # picture
                        img_count += 1
                        blob = shape.image.blob
                        mime = shape.image.content_type or "image/png"
                        vis = vision.analyze_image(prompt=prompt, image_bytes=blob, mime=mime)
                        if vis:
                            parts.append(f"\n--- SLIDE {si} IMAGE {img_count} VISION ---\n{vis}")
                except Exception:
                    continue

    return "\n".join(parts).strip()